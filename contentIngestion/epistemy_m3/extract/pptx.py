"""PPTX extractor (T9): one unit per slide with title, body, speaker notes."""
from __future__ import annotations
import io
from typing import List

from epistemy_m3.models import ChunkPosition
from epistemy_m3.chunking import ExtractedUnit


class PptxExtractor:
    """python-pptx per-slide extraction; empty slides flagged for the chunker."""

    def extract(self, data: bytes) -> List[ExtractedUnit]:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        return [self._slide_unit(s, i) for i, s in enumerate(prs.slides, start=1)]

    def _slide_unit(self, slide, slide_no: int) -> ExtractedUnit:
        """Combine title, body text, and notes into one structural unit."""
        title = self._title(slide)
        body = self._body_text(slide)
        notes = self._notes(slide)
        text = "\n".join(p for p in [body, notes] if p).strip()
        path = [title] if title else None
        return ExtractedUnit(text=text,
                             position=ChunkPosition(slide_no=slide_no, heading_path=path),
                             structured=bool(text))

    def _title(self, slide) -> str:
        """Return the slide title text, or empty when absent."""
        title_ph = getattr(slide.shapes, "title", None)
        return title_ph.text.strip() if title_ph and title_ph.text else ""

    def _body_text(self, slide) -> str:
        """Join text from every non-title shape that carries a text frame."""
        title_ph = getattr(slide.shapes, "title", None)
        parts = [s.text.strip() for s in slide.shapes
                 if s is not title_ph and s.has_text_frame and s.text.strip()]
        return "\n".join(parts)

    def _notes(self, slide) -> str:
        """Append speaker notes under a Notes separator when present."""
        if not slide.has_notes_slide:
            return ""
        note = slide.notes_slide.notes_text_frame.text.strip()
        return f"Notes: {note}" if note else ""

"""T9: per-slide units with slide_no and speaker notes."""
import io

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches

from epistemy_m3.extract.pptx import PptxExtractor


def _build_deck(n: int) -> bytes:
    """Create an n-slide deck with titles, body text, and notes."""
    prs = Presentation()
    for i in range(1, n + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Topic {i}"
        slide.placeholders[1].text = f"body of slide {i}"
        slide.notes_slide.notes_text_frame.text = f"note {i}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_one_unit_per_slide_with_numbers_and_notes():
    units = PptxExtractor().extract(_build_deck(30))
    assert len(units) == 30
    assert [u.position.slide_no for u in units] == list(range(1, 31))
    assert "Notes: note 1" in units[0].text
    assert units[0].position.heading_path == ["Topic 1"]
